
from tkinter import messagebox
import tkinter as tk
from tkinter import Tk, PhotoImage
from tkinter import filedialog, Text, messagebox
from docx import Document
import pandas as pd
from pptx import Presentation
from tkinter import filedialog, Text, messagebox
from tkinter.font import Font
from PIL import Image, ImageTk
from tkinter import font

def open_word_file():
    file_path = filedialog.askopenfilename(filetypes=[("Word Documents", "*.docx")])
    if file_path:
        doc = Document(file_path)
        content = ""
        for paragraph in doc.paragraphs:
            content += paragraph.text + "\n"
        text_widget.delete(1.0, tk.END)
        text_widget.insert(tk.END, content)
        status_bar.config(text=f"Opened Word file: {file_path}")

def save_word_file():
    content = text_widget.get(1.0, tk.END)
    save_path = filedialog.asksaveasfilename(defaultextension=".docx")
    if save_path:
        doc = Document()
        for line in content.split("\n"):
            doc.add_paragraph(line)
        doc.save(save_path)
        status_bar.config(text=f"Saved Word file: {save_path}")

def save_as_word_file():
    content = text_widget.get(1.0, tk.END)
    save_path = filedialog.asksaveasfilename(defaultextension=".docx")
    if save_path:
        doc = Document()
        for line in content.split("\n"):
            doc.add_paragraph(line)
        doc.save(save_path)
        status_bar.config(text=f"Saved Word file as: {save_path}")


# print function emplamented 💋💋



def open_excel_file():
    file_path = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx")])
    if file_path:
        df = pd.read_excel(file_path)
        content = df.to_string(index=False)
        text_widget.delete(1.0, tk.END)
        text_widget.insert(tk.END, content)
        status_bar.config(text=f"Opened Excel file: {file_path}")

def save_excel_file():
    content = text_widget.get(1.0, tk.END)
    save_path = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel Files", "*.xlsx")])
    if save_path:
        # Content ko DataFrame mein convert karen
        data = [line.split("\t") for line in content.strip().split("\n")]
        df = pd.DataFrame(data)
        # Excel file save karen
        with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
            df.to_excel(writer, index=False, header=False)
        status_bar.config(text=f"Saved Excel file: {save_path}")

def save_as_excel_file():
    content = text_widget.get(1.0, tk.END)
    save_path = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel Files", "*.xlsx")])
    if save_path:
        # Content ko DataFrame mein convert karen
        data = [line.split("\t") for line in content.strip().split("\n")]
        df = pd.DataFrame(data)
        # Excel file save karen
        with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
            df.to_excel(writer, index=False, header=False)
        status_bar.config(text=f"Saved Excel file as: {save_path}")

# new function emaplimet

def open_excel_file():
    file_path = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx")])
    if file_path:
        df = pd.read_excel(file_path)
        # DataFrame ko string mein convert karte waqt headers aur indexes include karen
        content = df.to_string(index=True, header=True)
        text_widget.delete(1.0, tk.END)
        text_widget.insert(tk.END, content)
        status_bar.config(text=f"Opened Excel file: {file_path}")


# align me function 

def set_text_center_aligned(text_widget):
    content = text_widget.get(1.0, tk.END)
    lines = content.strip().split("\n")
    max_width = max(len(line) for line in lines)
    centered_content = "\n".join(f"{line.center(max_width)}" for line in lines)
    text_widget.delete(1.0, tk.END)
    text_widget.insert(tk.END, centered_content)



def open_ppt_file():
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Presentations", "*.pptx")])
    if file_path:
        prs = Presentation(file_path)
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        text_widget.delete(1.0, tk.END)
        text_widget.insert(tk.END, content)
        status_bar.config(text=f"Opened PowerPoint file: {file_path}")

def copy_text():
    root.clipboard_clear()
    root.clipboard_append(text_widget.selection_get())
    status_bar.config(text="Copied text")

def cut_text():
    copy_text()
    text_widget.delete(tk.SEL_FIRST, tk.SEL_LAST)
    status_bar.config(text="Cut text")

def paste_text():
    text_widget.insert(tk.INSERT, root.clipboard_get())
    status_bar.config(text="Pasted text")

def about_app():
    messagebox.showinfo("About", "This is a simple file opener  application using .")

# Main window setup
root = tk.Tk()
bold_font = font.Font(family="Arial", size=12, weight="bold")
root.title("")
label = tk.Label(root, bg="white", text="Welcome to File Opener!", font=bold_font)
label.pack()
icon = PhotoImage(file="c:/Users/admin/Desktop/Amit-main/template/iconpoin.png")
root.wm_iconphoto(True, icon)

# Toolbar setup
toolbar = tk.Frame(root, bg='white', relief='raised', borderwidth=2)
toolbar.pack(side='top', fill='x')  # 'top' ensures it appears above the menu bar

img = Image.open("c:/Users/admin/Desktop/Amit-main/template/word.png")
icon_image = ImageTk.PhotoImage(img)
img_resized = img.resize((38, 38))
img_resized.save("c:/Users/admin/Desktop/Amit-main/template/word.png")
open_word_toolbar_btn = tk.Button(toolbar, bg="white", text="word",image=icon_image,command=open_word_file)
open_word_toolbar_btn.pack(side='left', padx=3,  pady=3)
open_word_toolbar_btn.image_names = icon_image

oon = Image.open("c:/Users/admin/Desktop/Amit-main/template/excel.png")
image_icon = ImageTk.PhotoImage(oon)
oon_resized = oon.resize((38, 38))
oon_resized.save("c:/Users/admin/Desktop/Amit-main/template/excel.png")
open_excel_toolbar_btn = tk.Button(toolbar, text="excel", image=image_icon, command=open_excel_file )
open_excel_toolbar_btn.pack(side="left", padx=3, pady=3)
open_excel_toolbar_btn.image_names = image_icon

oop = Image.open("c:/Users/admin/Desktop/Amit-main/template/powerpoint.png")
photo_icon = ImageTk.PhotoImage(oop)
oop_resized = oop.resize((37, 37))
oop_resized.save("c:/Users/admin/Desktop/Amit-main/template/powerpoint.png")
open_ppt_toolbar_btn = tk.Button(toolbar, text="powerpoint", image=photo_icon, command=open_ppt_file)
open_ppt_toolbar_btn.pack(side="left", padx=3, pady=3)
open_ppt_toolbar_btn.image_names = photo_icon




# Menu setup
menu_bar = tk.Menu(root)
root.config(menu=menu_bar)

# File Menu
file_menu = tk.Menu(menu_bar, tearoff=0)
menu_bar.add_cascade(label="File", menu=file_menu)
file_menu.add_command(label="Save Word File", command=save_word_file, accelerator="Ctrl+S")
file_menu.add_command(label="Save Word File As", command=save_as_word_file, accelerator="Ctrl+Shift+S")
file_menu.add_command(label="Save Excel File", command=save_excel_file, accelerator="Ctrl+E")
file_menu.add_command(label="Save Excel File As", command=save_as_excel_file, accelerator="Ctrl+Shift+E")
file_menu.add_separator()
file_menu.add_command(label="Open Word File", command=open_word_file)
file_menu.add_separator()
file_menu.add_command(label="Open Excel File", command=open_excel_file)
file_menu.add_separator()
file_menu.add_command(label="Open PowerPoint File", command=open_ppt_file)
file_menu.add_separator()
file_menu.add_command(label="Exit", command=root.quit)

home_menu = tk.Menu(menu_bar, tearoff=0)
home_menu.add_cascade(label="Home", menu=home_menu)
home_menu.add_command(label="nutrition")

# Edit Menu
edit_menu = tk.Menu(menu_bar, tearoff=0)
menu_bar.add_cascade(label="Edit", menu=edit_menu)
edit_menu.add_command(label="Cut", command=cut_text, accelerator="ftrl + x")
edit_menu.add_command(label="Copy", command=copy_text, accelerator="gtrl + c")
edit_menu.add_command(label="Paste", command=paste_text, accelerator="vtrl + v")
edit_menu.add_command(label="find", accelerator="ltrl + f" )
edit_menu.add_command(label="replace", accelerator="ctrl + r")
# Help Menu
help_menu = tk.Menu(menu_bar, tearoff=0)
menu_bar.add_cascade(label="Help", menu=help_menu)
help_menu.add_command(label="About", command=about_app)

#height
# Text box and scrollbar setup
text_frame = tk.Frame(root)
text_frame.pack(expand=True, fill='both', side='right')

vertical_scrollbar = tk.Scrollbar(text_frame)
vertical_scrollbar.pack(side='right', fill='y')

horizontal_scrollbar = tk.Scrollbar(text_frame, orient="horizontal")
horizontal_scrollbar.pack(side='bottom', fill='x')

text_widget = Text(text_frame, wrap="none", yscrollcommand=vertical_scrollbar.set, xscrollcommand=horizontal_scrollbar.set)
text_widget.pack(expand=True, fill='both')

vertical_scrollbar.config(command=text_widget.yview)
horizontal_scrollbar.config(command=text_widget.xview)

# Status bar setup
status_bar = tk.Label(root, bg="white", text="Ready", bd=1, relief="sunken", anchor="w")
status_bar.pack(side="bottom", fill="x")

class WordZoomApp:
    def __init__(self, root):
        self.root = root
        
        # Text Widget
        self.text_widget = tk.Text(root, wrap="word", font=("Arial", 12))
        self.text_widget.pack(expand=True, fill="both")

        # Zoom Slider
        self.zoom_slider = tk.Scale(root, bg="white", from_=1, to=30, orient="horizontal", command=self.adjust_zoom)
        self.zoom_slider.set(12)  # Default font size
        self.zoom_slider.pack(fill="x")

    def adjust_zoom(self, value):
        """Update font size based on slider value."""
        new_font = ("Arial", int(value))
        self.text_widget.config(font=new_font)


app = WordZoomApp(root)

# llll


# Shortcut keys setup
root.bind_all("<Control-o>", lambda event: open_word_file())
root.bind_all("<Control-s>", lambda event: save_word_file())
root.bind_all("<Control-Shift-S>", lambda event: save_as_word_file())
root.bind_all("<Control-e>", lambda event: save_excel_file())
root.bind_all("<Control-Shift-E>", lambda event: save_as_excel_file())


# new function

def save_excel_file():
    content = text_widget.get(1.0, tk.END)
    save_path = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel Files", "*.xlsx")])
    if save_path:
        # Content को DataFrame में बदलें
        data = [line.split("\t") for line in content.strip().split("\n")]
        df = pd.DataFrame(data)
        
        # Data को टाइप कास्ट करें (यदि जरूरी हो)
        try:
            numeric_df = df.apply(pd.to_numeric, errors='coerce')  # केवल संख्यात्मक कॉलम
            numeric_df.fillna(0, inplace=True)  # Missing values को 0 से भरें
            
            # गणना (Sum, Multiplication, Average)
            sum_row = numeric_df.sum(axis=0)
            product_row = numeric_df.prod(axis=0)
            avg_row = numeric_df.mean(axis=0)

            # अंतिम DataFrame में गणना जोड़ें
            numeric_df.loc['Sum'] = sum_row
            numeric_df.loc['Product'] = product_row
            numeric_df.loc['Average'] = avg_row
        except Exception as e:
            messagebox.showerror("Error", f"Calculation failed: {e}")
            return
        
        # Excel फाइल को सेव करें
        with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
            numeric_df.to_excel(writer, index=True, header=True)
        
        status_bar.config(text=f"Saved Excel file with calculations: {save_path}")



# Text box setup with the custom font
text_frame = tk.Frame(root)
text_frame.pack(expand=True, fill='both', side='right')

vertical_scrollbar = tk.Scrollbar(text_frame)
vertical_scrollbar.pack(side='right', fill='y')

horizontal_scrollbar = tk.Scrollbar(text_frame, orient="horizontal")
horizontal_scrollbar.pack(side='bottom', fill='x')


def open_find_replace_dialog(event=None):
    # Create the Find and Replace dialog
    dialog = tk.Toplevel(root, bg="lightblue" )
    dialog.title("Find and Replace")
    dialog.geometry("400x200")
    
    # Find Label and Entry
    tk.Label(dialog, text=" Find : ").grid(row=0, column=0, padx=20, pady=20)
    find_entry = tk.Entry(dialog, width=40)
    find_entry.grid(row=0, column=1, padx=20, pady=20)
    
    # Replace Label and Entry
    tk.Label(dialog, text=" Replace : ").grid(row=1, column=0, padx=20, pady=20)
    replace_entry = tk.Entry(dialog, width=40)
    replace_entry.grid(row=1, column=1, padx=20, pady=20)
    
    # Buttons for Find and Replace
    def perform_find():
        query = find_entry.get()
        text_widget.tag_remove("highlights", "1.0", tk.END)  # Remove previous highlights
        if query:
            start_pos = "1.0"
            while True:
                start_pos = text_widget.search(query, start_pos, stopindex=tk.END)
                if not start_pos:
                    break
                end_pos = f"{start_pos}+{len(query)}c"
                text_widget.tag_add("highlights", start_pos, end_pos)
                start_pos = end_pos
            text_widget.tag_config("highlights", background="yellow")
            status_bar.config(text=f"Search completed for: {query}")
    
    def perform_replace():
        find_text = find_entry.get()
        replace_text = replace_entry.get()
        content = text_widget.get(1.0, tk.END)
        updated_content = content.replace(find_text, replace_text)
        text_widget.delete(1.0, tk.END)
        text_widget.insert(tk.END, updated_content)
        status_bar.config(text=f"Replaced '{find_text}' with '{replace_text}'")
    
    tk.Button(dialog, text="Find", command=perform_find).grid(row=2, column=0, padx=15, pady=15)
    find_button = tk.Button(dialog, text="find", command=perform_find, width=10 )
    find_button.grid(row=2, column=0, padx=15, pady=15)
    tk.Button(dialog, text="Replace", command=perform_replace).grid(row=2, column=1, padx=15, pady=15)
    replace_button = tk.Button(dialog, text="replace", command=perform_replace, width=12)
    replace_button.grid(row=2, column=1, padx=15, pady=15 )
    # Make sure the dialog doesn't block the main window
    dialog.transient(root)
    dialog.grab_set()
    dialog.focus_set()
    dialog.wait_window(dialog)

# Bind the dialog to Shortcut Keys
root.bind("<Control-f>", open_find_replace_dialog)  # Ctrl+F to open Find dialog
root.bind("<Control-r>", open_find_replace_dialog)  # Ctrl+R to open Replace dialog


# main loop
root.mainloop()
